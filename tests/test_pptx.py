from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.converters.pptx import PptxConverter


def _make_pptx(path: Path) -> None:
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Título da Apresentação"
    slide1.placeholders[1].text = "Introdução ao projeto"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body = slide2.placeholders[1]
    text_frame = body.text_frame
    text_frame.text = "Primeiro item"
    paragraph = text_frame.add_paragraph()
    paragraph.text = "Segundo item"

    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Conclusão"
    prs.save(str(path))


def test_convert_pptx(tmp_path: Path) -> None:
    source = tmp_path / "apresentacao.pptx"
    _make_pptx(source)

    result = PptxConverter().convert(source, tmp_path / "output")

    assert result.markdown
    assert "Título da Apresentação" in result.markdown
    assert "Introdução ao projeto" in result.markdown
    assert "Agenda" in result.markdown
    assert "Primeiro item" in result.markdown
    assert "Conclusão" in result.markdown
